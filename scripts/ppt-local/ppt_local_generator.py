#!/usr/bin/env python3
"""
ppt_local_generator.py — P0 本地 LLM + python-pptx PPT 生成器
===============================================================
用法:
  python3 ppt_local_generator.py <json_payload_file>
  
输入 JSON 格式:
{
  "title": "PPT标题",
  "slides": [
    {
      "title": "页面标题",
      "content": ["要点1", "要点2", "要点3"],
      "type": "bullets" | "title" | "toc" | "two-column" | "chart" | "quote"
    }
  ],
  "theme": "consultant" | "founder" | "blues" | ...,
  "tone": "professional" | "casual" | "creative",
  "output_path": "/tmp/xxx.pptx"
}

依赖:
  pip install python-pptx lxml
"""

import json
import sys
import os
import subprocess
import re
from pathlib import Path

# 尝试导入 python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print(json.dumps({"error": "python-pptx 未安装，请运行: pip install python-pptx"}))
    sys.exit(1)


# ==================== 主题配色数据库 ====================

THEME_COLORS = {
    # theme_id: (primary, secondary, accent, background, text)
    "consultant": {
        "primary": RGBColor(0x1A, 0x3A, 0x5C),   # 深蓝
        "secondary": RGBColor(0x2C, 0x5A, 0x8C), # 中蓝
        "accent": RGBColor(0xD4, 0xAF, 0x37),     # 金色
        "background": RGBColor(0xFF, 0xFF, 0xFF), # 白色
        "text": RGBColor(0x1F, 0x29, 0x37),       # 深灰
        "subtitle": RGBColor(0x6B, 0x72, 0x80),   # 浅灰
    },
    "founder": {
        "primary": RGBColor(0x1E, 0x3A, 0x5F),
        "secondary": RGBColor(0x2E, 0x5A, 0x8F),
        "accent": RGBColor(0xE8, 0xB9, 0x4B),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "blues": {
        "primary": RGBColor(0x17, 0x3A, 0x8C),
        "secondary": RGBColor(0x1E, 0x5A, 0x9C),
        "accent": RGBColor(0x00, 0x96, 0x88),
        "background": RGBColor(0xF0, 0xF4, 0xF8),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "electric": {
        "primary": RGBColor(0x7C, 0x3A, 0xED),
        "secondary": RGBColor(0xA8, 0x55, 0xF2),
        "accent": RGBColor(0xF5, 0x72, 0x3A),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "icebreaker": {
        "primary": RGBColor(0x06, 0xB6, 0xD4),
        "secondary": RGBColor(0x38, 0xB2, 0xAC),
        "accent": RGBColor(0xED, 0x89, 0x3A),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "aurora": {
        "primary": RGBColor(0x06, 0xB6, 0xD4),
        "secondary": RGBColor(0x8B, 0x5C, 0xF6),
        "accent": RGBColor(0x10, 0xB9, 0x81),
        "background": RGBColor(0x0F, 0x17, 0x29),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle": RGBColor(0x9C, 0xA3, 0xAF),
    },
    "ash": {
        "primary": RGBColor(0x37, 0x4A, 0x5C),
        "secondary": RGBColor(0x5C, 0x6F, 0x7D),
        "accent": RGBColor(0x94, 0xA3, 0xB8),
        "background": RGBColor(0xF8, 0xF9, 0xFA),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "default": {
        "primary": RGBColor(0x25, 0x63, 0xEB),
        "secondary": RGBColor(0x37, 0x4A, 0x5C),
        "accent": RGBColor(0x7C, 0x3A, 0xED),
        "background": RGBColor(0xFF, 0xFF, 0xFF),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
    "festival": {
        "primary": RGBColor(0xC0, 0x23, 0x2C),
        "secondary": RGBColor(0x1A, 0x3A, 0x5C),
        "accent": RGBColor(0xD4, 0xAF, 0x37),
        "background": RGBColor(0xFF, 0xF8, 0xF0),
        "text": RGBColor(0x1F, 0x29, 0x37),
        "subtitle": RGBColor(0x6B, 0x72, 0x80),
    },
}

DEFAULT_THEME = THEME_COLORS["default"]


# ==================== LLM 调用 ====================

def call_llm(messages, api_base=None, api_key=None, model="MiniMax-M2", timeout=60):
    """
    调用本地 LLM (MiniMax / GLM 兼容 OpenAI 格式)
    
    messages: [{"role": "user"|"system"|"assistant", "content": "..."}]
    """
    import urllib.request
    import urllib.error
    
    if not api_base:
        api_base = os.environ.get("GLM_API_BASE", "https://mydamoxing.cn/v1/chat/completions")
    if not api_key:
        api_key = os.environ.get("GLM_API_KEYS", os.environ.get("MYDAMOXING_API_KEY", ""))
    
    if not api_key:
        return None, "LLM API Key 未配置"
    
    payload = {
        "model": model,
        "messages": messages,
        "max_tokens": 8192,
        "temperature": 0.7,
    }
    
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        api_base,
        data=data,
        headers={
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        },
        method="POST"
    )
    
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            result = json.loads(resp.read().decode("utf-8"))
            content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
            return content, None
    except urllib.error.HTTPError as e:
        err_body = e.read().decode("utf-8") if e.fp else ""
        return None, f"HTTP {e.code}: {err_body[:200]}"
    except Exception as e:
        return None, str(e)


# ==================== PPT 生成核心 ====================

def get_theme(theme_id):
    """获取主题配色"""
    return THEME_COLORS.get(theme_id, DEFAULT_THEME)


def add_title_slide(prs, title, subtitle, theme):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 全幅背景色
    background = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = theme["primary"]
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.8)
        )
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.alignment = PP_ALIGN.CENTER
    
    # 底部装饰条
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(6.8), prs.slide_width, Inches(0.2)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["accent"]
    accent_bar.line.fill.background()
    
    return slide


def add_content_slide(prs, page_title, bullets, theme):
    """添加内容页（要点列表）"""
    slide_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["background"]
    bg.line.fill.background()
    
    # 左侧色条装饰
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.15), prs.slide_height
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["primary"]
    accent_bar.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = page_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]
    
    # 分隔线
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.1), Inches(11.5), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme["accent"]
    line.line.fill.background()
    
    # 要点列表
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4), Inches(12), Inches(4.5)
    )
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = theme["text"]
        p.space_before = Pt(12)
        p.space_after = Pt(4)
    
    return slide


def add_two_column_slide(prs, page_title, left_content, right_content, theme):
    """添加双栏页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["background"]
    bg.line.fill.background()
    
    # 页面标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = page_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["primary"]
    
    # 左栏
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(5.8), Inches(4.8)
    )
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    for i, item in enumerate(left_content):
        p = tf_l.paragraphs[i] if i == 0 else tf_l.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = theme["text"]
        p.space_before = Pt(8)
    
    # 右栏
    right_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.2), Inches(5.8), Inches(4.8)
    )
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    for i, item in enumerate(right_content):
        p = tf_r.paragraphs[i] if i == 0 else tf_r.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = theme["text"]
        p.space_before = Pt(8)
    
    # 中间分隔线
    sep_line = slide.shapes.add_shape(
        1, Inches(6.4), Inches(1.2), Pt(2), Inches(4.5)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = theme["accent"]
    sep_line.line.fill.background()
    
    return slide


def add_ending_slide(prs, title, theme):
    """添加结尾页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["primary"]
    bg.line.fill.background()
    
    # 装饰线
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(3.0), prs.slide_width, Pt(3)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["accent"]
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"谢谢观看"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.6)
    )
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    
    return slide


def generate_pptx(data):
    """
    根据输入数据生成 PPTX
    data: {
        "title": str,
        "slides": [{"title": str, "content": [str], "type": str}],
        "theme": str,
        "tone": str,
        "output_path": str
    }
    """
    title = data.get("title", "PPT")
    slides = data.get("slides", [])
    theme_id = data.get("theme", "default")
    output_path = data.get("output_path", "/tmp/output.pptx")
    
    theme = get_theme(theme_id)
    
    # 创建演示文稿（16:9）
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 封面
    add_title_slide(prs, title, None, theme)
    
    # 内容页
    for idx, slide_data in enumerate(slides):
        slide_title = slide_data.get("title", f"第{idx+1}页")
        content = slide_data.get("content", [])
        slide_type = slide_data.get("type", "bullets")
        
        if slide_type == "two-column" and len(content) >= 2:
            mid = len(content) // 2
            add_two_column_slide(
                prs, slide_title, content[:mid], content[mid:], theme
            )
        else:
            add_content_slide(prs, slide_title, content, theme)
    
    # 结尾页
    add_ending_slide(prs, title, theme)
    
    # 保存
    prs.save(output_path)
    return output_path


# ==================== LLM 增强内容（可选） ====================

def enhance_with_llm(raw_input, theme_id, tone, api_base=None, api_key=None):
    """
    调用本地 LLM 将用户输入增强为结构化 PPT 数据
    raw_input: 原始文本输入
    返回: {"title": str, "slides": [...]}
    """
    system_prompt = f"""你是一个专业的PPT内容规划助手。用户会提供一段文字内容，你需要将其转化为PPT大纲。

要求：
1. 分析内容，提炼出5-8个核心要点
2. 每个页面只放3-4个要点
3. 返回JSON格式的大纲，包含title和slides数组
4. 每页要有title、content（要点数组）、type（bullets/two-column等）

输出格式（只返回JSON，不要其他内容）：
{{
  "title": "PPT标题",
  "slides": [
    {{"title": "页面标题", "content": ["要点1", "要点2", "要点3"], "type": "bullets"}},
    ...
  ]
}}

主题色系: {theme_id}
风格: {tone}
"""
    
    user_prompt = f"""请将以下内容转化为PPT大纲：

{raw_input}
"""
    
    content, err = call_llm(
        [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        api_base=api_base,
        api_key=api_key,
    )
    
    if err:
        return None, err
    
    # 尝试提取 JSON
    try:
        # 尝试找到 ```json ... ``` 或 ``` ... ``` 包裹的 JSON
        json_match = re.search(r'```(?:json)?\s*([\s\S]+?)\s*```', content)
        if json_match:
            json_str = json_match.group(1)
        else:
            # 直接解析整段
            json_str = content.strip()
        
        result = json.loads(json_str)
        return result, None
    except json.JSONDecodeError as e:
        return None, f"JSON解析失败: {e}\n原始内容: {content[:300]}"


# ==================== 主入口 ====================

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "用法: python3 ppt_local_generator.py <json_payload_file>"}))
        sys.exit(1)
    
    payload_file = sys.argv[1]
    
    # 读取输入
    if payload_file == "-":
        # 从 stdin 读取
        raw_input = sys.stdin.read()
    else:
        if not os.path.exists(payload_file):
            print(json.dumps({"error": f"文件不存在: {payload_file}"}))
            sys.exit(1)
        with open(payload_file, "r", encoding="utf-8") as f:
            raw_input = f.read()
    
    try:
        data = json.loads(raw_input)
    except json.JSONDecodeError as e:
        print(json.dumps({"error": f"JSON解析失败: {e}"}))
        sys.exit(1)
    
    # 获取 API 配置
    api_base = data.pop("api_base", None) or os.environ.get("GLM_API_BASE")
    api_key = data.pop("api_key", None) or os.environ.get("GLM_API_KEYS") or os.environ.get("MYDAMOXING_API_KEY")
    use_llm = data.pop("use_llm", False)
    
    # 如果需要 LLM 增强
    if use_llm and data.get("raw_input"):
        enhanced, err = enhance_with_llm(
            data["raw_input"],
            data.get("theme", "default"),
            data.get("tone", "professional"),
            api_base=api_base,
            api_key=api_key,
        )
        if err:
            print(json.dumps({"error": f"LLM 增强失败: {err}"}))
            sys.exit(1)
        data = enhanced
    
    # 生成 PPTX
    output_path = data.get("output_path") or f"/tmp/ppt_{os.getpid()}.pptx"
    data["output_path"] = output_path
    
    try:
        result_path = generate_pptx(data)
        result = {
            "success": True,
            "output_path": result_path,
            "file_size": os.path.getsize(result_path),
        }
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({"error": f"PPTX 生成失败: {e}"}))
        sys.exit(1)


if __name__ == "__main__":
    main()
