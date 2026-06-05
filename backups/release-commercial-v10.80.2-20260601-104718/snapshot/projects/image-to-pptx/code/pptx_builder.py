"""
PPTX Builder - PPTX 构建器

职责：
    使用 python-pptx 生成真实 PPTX 文件

核心逻辑：
    - 16:9比例（13.333 x 7.5 英寸）
    - 布局解析：将设计图的布局信息转换为 python-pptx 坐标
    - 配色应用：将设计图的配色方案应用到 PPT 的背景、文本、形状

技术栈：
    - python-pptx >= 0.6.21
    - Pillow >= 10.0.0
"""

import os
from dataclasses import dataclass
from pathlib import Path
from typing import Optional
from PIL import Image

# python-pptx 相关导入
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


@dataclass
class LayoutConfig:
    """布局配置"""
    # 幻灯片尺寸（16:9）
    slide_width: float = 13.333  # 英寸
    slide_height: float = 7.5    # 英寸
    
    # 边距
    margin_left: float = 0.5
    margin_right: float = 0.5
    margin_top: float = 0.4
    margin_bottom: float = 0.4
    
    # 标题区域
    title_top: float = 0.3
    title_height: float = 0.8
    
    # 内容区域
    content_top: float = 1.2
    content_height: float = 5.5
    
    # 页脚区域
    footer_top: float = 7.0
    footer_height: float = 0.4
    
    # 页码位置
    page_number_right: float = 0.5


@dataclass
class ColorScheme:
    """配色方案"""
    primary: str = "#2563EB"
    secondary: str = "#3B82F6"
    background: str = "#FFFFFF"
    accent: str = "#06B6D4"
    text_dark: str = "#1F2937"
    text_light: str = "#6B7280"
    
    def to_pptx_colors(self) -> dict:
        """转换为 python-pptx 颜色格式"""
        return {
            "primary": self._hex_to_rgb(self.primary),
            "secondary": self._hex_to_rgb(self.secondary),
            "background": self._hex_to_rgb(self.background),
            "accent": self._hex_to_rgb(self.accent),
            "text_dark": self._hex_to_rgb(self.text_dark),
            "text_light": self._hex_to_rgb(self.text_light),
        }
    
    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """HEX 颜色转换为 RGBColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)


@dataclass
class SlideContent:
    """幻灯片内容"""
    title: str
    content: Optional[list[str]] = None
    bullet_points: Optional[list[str]] = None
    image_path: Optional[str] = None
    background_color: Optional[str] = None


class PPTXBuilder:
    """
    PPTX 构建器
    
    使用 python-pptx 生成专业 PPTX 文件。
    支持自定义布局、配色、样式。
    """
    
    def __init__(
        self,
        output_path: str = "./output.pptx",
        layout_config: Optional[LayoutConfig] = None,
        color_scheme: Optional[ColorScheme] = None,
    ):
        """
        初始化 PPTX 构建器
        
        Args:
            output_path: 输出文件路径
            layout_config: 布局配置
            color_scheme: 配色方案
        """
        self.output_path = Path(output_path)
        self.layout = layout_config or LayoutConfig()
        self.colors = color_scheme or ColorScheme()
        self.prs = None
        self._slide_count = 0
    
    def create_presentation(self) -> Presentation:
        """创建演示文稿"""
        self.prs = Presentation()
        # 设置 16:9 比例
        self.prs.slide_width = Inches(self.layout.slide_width)
        self.prs.slide_height = Inches(self.layout.slide_height)
        self._slide_count = 0
        return self.prs
    
    def add_slide(
        self,
        slide_content: SlideContent,
        layout_index: int = 6,  # 6 是空白布局
    ) -> Presentation:
        """
        添加幻灯片
        
        Args:
            slide_content: 幻灯片内容
            layout_index: 布局索引（默认空白布局）
            
        Returns:
            添加的幻灯片
        """
        if not self.prs:
            self.create_presentation()
        
        # 添加幻灯片
        slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[layout_index]
        )
        self._slide_count += 1
        
        # 设置背景色
        if slide_content.background_color:
            self._set_slide_background(slide, slide_content.background_color)
        
        # 添加标题
        if slide_content.title:
            self._add_title(slide, slide_content.title)
        
        # 添加内容
        if slide_content.bullet_points:
            self._add_bullet_content(
                slide, slide_content.bullet_points
            )
        elif slide_content.content:
            self._add_text_content(slide, slide_content.content)
        
        # 添加图片
        if slide_content.image_path:
            self._add_image(
                slide, slide_content.image_path
            )
        
        # 添加页码
        self._add_page_number(slide, self._slide_count)
        
        return slide
    
    def _set_slide_background(self, slide, color_hex: str) -> None:
        """设置幻灯片背景色"""
        from pptx.oxml.ns import qn
        
        hex_color = color_hex.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    def _add_title(self, slide, title: str) -> None:
        """添加标题"""
        left = Inches(self.layout.margin_left)
        top = Inches(self.layout.title_top)
        width = Inches(self.layout.slide_width - self.layout.margin_left - self.layout.margin_right)
        height = Inches(self.layout.title_height)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.text = title
        
        # 设置标题样式
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors.to_pptx_colors()["text_dark"]
        paragraph.alignment = PP_ALIGN.LEFT
        
        # 标题下方添加装饰线
        self._add_title_decoration(slide, top, height)
    
    def _add_title_decoration(self, slide, title_top: float, title_height: float) -> None:
        """添加标题装饰线"""
        left = Inches(self.layout.margin_left)
        top = Inches(title_top + title_height + 0.1)
        width = Inches(1.5)
        height = Inches(0.05)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors.to_pptx_colors()["primary"]
        shape.line.fill.background()  # 无边框
    
    def _add_bullet_content(self, slide, bullet_points: list[str]) -> None:
        """添加要点内容"""
        left = Inches(self.layout.margin_left)
        top = Inches(self.layout.content_top)
        width = Inches(
            self.layout.slide_width - self.layout.margin_left - self.layout.margin_right
        )
        height = Inches(self.layout.content_height)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors.to_pptx_colors()["text_dark"]
            p.space_after = Pt(12)
    
    def _add_text_content(self, slide, content: list[str]) -> None:
        """添加文本内容"""
        left = Inches(self.layout.margin_left)
        top = Inches(self.layout.content_top)
        width = Inches(
            self.layout.slide_width - self.layout.margin_left - self.layout.margin_right
        )
        height = Inches(self.layout.content_height)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = text
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors.to_pptx_colors()["text_dark"]
            p.space_after = Pt(8)
    
    def _add_image(self, slide, image_path: str) -> None:
        """添加图片"""
        if not os.path.exists(image_path):
            print(f"警告: 图片不存在 {image_path}")
            return
        
        # 计算图片位置和大小
        img = Image.open(image_path)
        img_width, img_height = img.size
        
        # 保持 16:9 比例
        available_width = self.layout.slide_width - self.layout.margin_left - self.layout.margin_right
        available_height = self.layout.content_height
        
        # 计算缩放比例
        scale_w = available_width / (img_width / 96)  # 假设 96 DPI
        scale_h = available_height / (img_height / 96)
        scale = min(scale_w, scale_h, 1.0)
        
        width = Inches((img_width / 96) * scale)
        height = Inches((img_height / 96) * scale)
        
        left = Inches(self.layout.margin_left + (available_width - width.inches) / 2)
        top = Inches(self.layout.content_top + (available_height - height.inches) / 2)
        
        slide.shapes.add_picture(image_path, left, top, width, height)
    
    def _add_page_number(self, slide, page_num: int) -> None:
        """添加页码"""
        left = Inches(
            self.layout.slide_width - self.layout.margin_right - 1
        )
        top = Inches(self.layout.footer_top)
        width = Inches(1)
        height = Inches(self.layout.footer_height)
        
        page_box = slide.shapes.add_textbox(left, top, width, height)
        tf = page_box.text_frame
        tf.text = f"{page_num}"
        
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = self.colors.to_pptx_colors()["text_light"]
        paragraph.alignment = PP_ALIGN.RIGHT
    
    def save(self, output_path: Optional[str] = None) -> str:
        """
        保存 PPTX 文件
        
        Args:
            output_path: 输出路径（可选，默认使用初始化时的路径）
            
        Returns:
            保存的文件路径
        """
        if not self.prs:
            raise RuntimeError("没有可保存的演示文稿，请先调用 add_slide")
        
        save_path = Path(output_path) if output_path else self.output_path
        save_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(save_path))
        return str(save_path)
    
    @classmethod
    def from_design_image(
        cls,
        design_image_path: str,
        output_path: str = "./output.pptx",
    ) -> "PPTXBuilder":
        """
        从设计图创建 PPTX 构建器
        
        Args:
            design_image_path: 设计图路径
            output_path: 输出路径
            
        Returns:
            配置好的 PPTXBuilder 实例
        """
        # TODO: 使用 AI 视觉分析设计图，提取布局和配色信息
        # 目前返回默认配置
        return cls(output_path=output_path)
    
    def create_title_slide(self, title: str, subtitle: Optional[str] = None) -> None:
        """创建标题页"""
        if not self.prs:
            self.create_presentation()
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(self.layout.margin_left),
            Inches(self.layout.slide_height / 2 - 0.5),
            Inches(self.layout.slide_width - self.layout.margin_left - self.layout.margin_right),
            Inches(1)
        )
        tf = title_box.text_frame
        tf.text = title
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.colors.to_pptx_colors()["text_dark"]
        paragraph.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(self.layout.margin_left),
                Inches(self.layout.slide_height / 2 + 0.5),
                Inches(self.layout.slide_width - self.layout.margin_left - self.layout.margin_right),
                Inches(0.5)
            )
            tf = subtitle_box.text_frame
            tf.text = subtitle
            paragraph = tf.paragraphs[0]
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = self.colors.to_pptx_colors()["text_light"]
            paragraph.alignment = PP_ALIGN.CENTER
    
    def create_section_slide(self, section_title: str) -> None:
        """创建章节分隔页"""
        if not self.prs:
            self.create_presentation()
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors.to_pptx_colors()["primary"]
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(self.layout.margin_left),
            Inches(self.layout.slide_height / 2 - 0.5),
            Inches(self.layout.slide_width - self.layout.margin_left - self.layout.margin_right),
            Inches(1)
        )
        tf = title_box.text_frame
        tf.text = section_title
        paragraph = tf.paragraphs[0]
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
