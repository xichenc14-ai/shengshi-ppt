"""
PPTX 构建器 — 主构建器
"""

import logging
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from pdf_to_pptx.converter_types import PageContent
from pdf_to_pptx.builder.text_builder import TextBuilder
from pdf_to_pptx.builder.shape_builder import ShapeBuilder
from pdf_to_pptx.builder.image_builder import ImageBuilder
from pdf_to_pptx.builder.table_builder import TableBuilder

logger = logging.getLogger("pdf_to_pptx")


class PPTXBuilder:
    """
    PPTX 文件主构建器

    使用 python-pptx 构建幻灯片
    """

    def __init__(
        self,
        resource_dir: str = "resources",
        options: Dict[str, Any] = None,
    ):
        """
        Args:
            resource_dir: 资源文件目录
            options: 构建选项
        """
        self.resource_dir = resource_dir
        self.options = options or {}
        self.prs = Presentation()

        # 设置幻灯片尺寸（默认 16:9）
        slide_width = self.options.get("slide_width_inch", 13.333)
        slide_height = self.options.get("slide_height_inch", 7.5)
        self.prs.slide_width = Inches(slide_width)
        self.prs.slide_height = Inches(slide_height)

        self.pages: List[PageContent] = []

    def build(self, pages: List[PageContent]) -> None:
        """根据 PageContent 列表构建所有幻灯片"""
        for page in pages:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
            self._build_slide(slide, page)

    def _build_slide(self, slide, page: PageContent) -> None:
        """构建单页幻灯片"""
        text_builder = TextBuilder(slide, page)
        shape_builder = ShapeBuilder(slide)
        image_builder = ImageBuilder(slide, self.resource_dir)
        table_builder = TableBuilder(slide)

        # 文本块
        for tb in page.text_blocks:
            text_builder.add_text(tb)

        # 图形块
        for sb in page.shape_blocks:
            shape_builder.add_shape(sb)

        # 图片块
        for ib in (page.images_extracted or []):
            image_builder.add_image(ib)

        # 表格块
        for tbl in page.table_blocks:
            table_builder.add_table(tbl)

    def save(self, output_path: str) -> None:
        """保存 PPTX 文件"""
        self.prs.save(output_path)
        logger.info(f"PPTX saved: {output_path}")